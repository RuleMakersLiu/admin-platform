"""文档文本抽取：把 PDF / Word / Excel / PPT / 纯文本附件归一化为纯文本。

抽取失败返回空串（由调用方决定如何提示），绝不抛异常打断对话。
旧版二进制 Office（.doc/.xls/.ppt，OLE 格式）不在支持范围 —— 建议另存为
.docx/.xlsx/.pptx 或转成 PDF / 图片上传。
"""
import io
import itertools
import logging
import os
from typing import Optional

logger = logging.getLogger(__name__)

# 纯文本类扩展名（直接 decode）
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json", ".yaml", ".yml",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".php", ".rb", ".rs",
    ".c", ".h", ".cpp", ".hpp", ".cs", ".swift", ".kt", ".scala", ".sh", ".sql",
    ".html", ".htm", ".css", ".xml", ".ini", ".toml", ".cfg", ".conf", ".log",
}

# mime -> kind；kind 决定用哪个抽取器
_MIME_KIND = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/csv": "csv",
}
_EXT_KIND = {
    ".pdf": "pdf", ".docx": "docx", ".xlsx": "xlsx", ".pptx": "pptx",
    ".xls": "xls",  # 旧版 BIFF Excel（xlrd 解析）
    ".csv": "csv",
    ".doc": "ole", ".ppt": "ole",  # .doc/.ppt 旧版仍不支持
}


def _kind(mime: str, filename: str) -> str:
    mime = (mime or "").split(";")[0].strip().lower()
    if mime in _MIME_KIND:
        return _MIME_KIND[mime]
    if mime.startswith("text/"):
        return "text"
    ext = os.path.splitext(filename or "")[1].lower()
    return _EXT_KIND.get(ext, "text" if ext in _TEXT_EXTS else "")


def extract_text_from_bytes(data: bytes, mime: str = "", filename: str = "") -> str:
    """从文档字节抽取纯文本。返回 '' 表示无文本或抽取失败。"""
    if not data:
        return ""
    kind = _kind(mime, filename)
    try:
        if kind == "pdf":
            return _extract_pdf(data)
        if kind == "docx":
            return _extract_docx(data)
        if kind == "xlsx":
            return _extract_xlsx(data)
        if kind == "xls":
            return _extract_xls(data)
        if kind == "csv":
            return _extract_csv(data)
        if kind == "pptx":
            return _extract_pptx(data)
        if kind == "ole":
            logger.info("legacy Office (.doc/.ppt) not supported: %s", filename)
            return ""
        # text 或未知：尝试 UTF-8 解码
        return data.decode("utf-8", errors="ignore")
    except Exception as exc:  # 抽取不能打断对话
        logger.warning("doc_extract failed for %s (%s): %s", filename, kind, exc)
        return ""


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            pages.append("")
    return "\n".join(pages).strip()


def _extract_docx(data: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts).strip()


def _rows_to_markdown_table(rows: list[list[str]]) -> str:
    """二维行（已去空行）→ markdown 表格（首行表头 + 分隔行）。列对齐补齐。"""
    if not rows:
        return ""
    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]
    lines = ["| " + " | ".join(rows[0]) + " |"]
    lines.append("|" + "|".join([" --- "] * ncols) + "|")
    lines.extend("| " + " | ".join(r) + " |" for r in rows[1:])
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            blocks.append(f"## {ws.title}\n" + _rows_to_markdown_table(rows))
    wb.close()
    return "\n\n".join(blocks).strip()


def _extract_xls(data: bytes) -> str:
    """旧版 .xls（BIFF）→ markdown 表格。xlrd 2.0+ 仅支持 .xls。"""
    import xlrd

    wb = xlrd.open_workbook(file_contents=data)
    blocks = []
    for ws in wb.sheets():
        rows = []
        for rx in range(ws.nrows):
            cells = [str(ws.cell_value(rx, cx)) for cx in range(ws.ncols)]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            blocks.append(f"## {ws.name}\n" + _rows_to_markdown_table(rows))
    return "\n\n".join(blocks).strip()


def _extract_csv(data: bytes) -> str:
    """CSV → markdown 表格（首行表头）。取前 200 行控体积。"""
    import csv

    text = data.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    rows = [[c for c in row] for row in itertools.islice(reader, 200) if any(c.strip() for c in row)]
    return _rows_to_markdown_table(rows) if rows else text.strip()


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append(f"## Slide {idx}\n" + "\n".join(texts))
    return "\n\n".join(slides).strip()
