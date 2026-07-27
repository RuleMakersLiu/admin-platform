"""Tests for doc_extract — builds synthetic docs offline (no network)."""
import io

from app.ai import doc_extract


def test_extract_text_plain():
    assert doc_extract.extract_text_from_bytes(b"hello world", "text/plain", "a.txt") == "hello world"


def test_extract_text_code_ext_falls_back_to_utf8():
    out = doc_extract.extract_text_from_bytes(b"print('hi')", "text/x-python", "a.py")
    assert "print" in out


def test_extract_docx_roundtrip():
    import docx

    d = docx.Document()
    d.add_paragraph("需求标题")
    d.add_paragraph("第二个段落")
    buf = io.BytesIO()
    d.save(buf)
    out = doc_extract.extract_text_from_bytes(
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "r.docx",
    )
    assert "需求标题" in out and "第二个段落" in out


def test_extract_xlsx_roundtrip():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["姓名", "年龄"])
    ws.append(["张三", "30"])
    buf = io.BytesIO()
    wb.save(buf)
    out = doc_extract.extract_text_from_bytes(
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "r.xlsx",
    )
    assert "张三" in out and "姓名" in out


def test_extract_pptx_roundtrip():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = "标题幻灯片"
    buf = io.BytesIO()
    prs.save(buf)
    out = doc_extract.extract_text_from_bytes(
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "r.pptx",
    )
    assert "标题幻灯片" in out


def test_extract_invalid_pdf_returns_empty_gracefully():
    assert doc_extract.extract_text_from_bytes(b"not a pdf", "application/pdf", "r.pdf") == ""


def test_extract_empty_bytes():
    assert doc_extract.extract_text_from_bytes(b"", "application/pdf", "r.pdf") == ""


def test_kind_dispatch():
    assert doc_extract._kind("application/pdf", "") == "pdf"
    assert doc_extract._kind("", "r.docx") == "docx"
    assert doc_extract._kind("", "d.xlsx") == "xlsx"
    assert doc_extract._kind("", "deck.pptx") == "pptx"
    assert doc_extract._kind("", "legacy.doc") == "ole"
    assert doc_extract._kind("text/csv", "") == "text"
